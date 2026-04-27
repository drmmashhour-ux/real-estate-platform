#!/usr/bin/env python3
"""Generate LECIPM pitch deck. Requires: pip install python-pptx"""

from pathlib import Path

from pptx import Presentation

# Output next to this script’s repo root (scripts/ is under repo)
OUT = Path(__file__).resolve().parent.parent / "LECIPM_pitch.pptx"

slides_content = [
    ("LECIPM", "AI-powered real estate marketplace"),
    ("Problem", "Marketplaces are static\nManual pricing\nLow conversion"),
    ("Solution", "Self-optimizing system\nAI pricing\nAI growth"),
    ("Product", "Search\nFeed\nPricing\nCampaigns\nTrust\nBooking"),
    ("How it works", "Data → AI → Optimization → Growth"),
    ("Demo", "Search → Pricing → Booking → Campaigns → AI Brain"),
    ("Why now", "AI + personalization + efficiency"),
    ("Business model", "Fees\nBroker tools\nGrowth SaaS"),
    ("Traction", "Live system\nReady to launch"),
    ("Vision", "Operating system for real estate\n10-min demo?"),
]

if __name__ == "__main__":
    prs = Presentation()
    for title, content in slides_content:
        prs.slides.add_slide(prs.slide_layouts[1])
        slide = prs.slides[-1]
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
    prs.save(OUT)
    print(f"Wrote {OUT}")
