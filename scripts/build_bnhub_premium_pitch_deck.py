#!/usr/bin/env python3
"""
Build BNHub_Pitch_Deck_Premium.pptx — black (#0B0B0B) + gold (#D4AF37) investor deck.

- If docs/presentations/BNHub_Pitch_Deck.pptx exists, copies it to archive/ before generating.
- After building the 16-slide template, **merges content verbatim** from BNHub_Pitch_Deck.pptx (text, pictures, tables, autoshapes) and reapplies **#0B0B0B** slide background. Default map skips source “Vision” (slide 14) so premium slide 14 = source “The opportunity” (15); override via `MERGE_TARGET_TO_SOURCE` in this file.
- Embeds product screenshots from docs/presentations/screenshots/*.png when merge is off or for template-only rebuild (`PITCH_MERGE=0`).

Run from repo root:
  cd docs/presentations && python3 -m venv .venv && source .venv/bin/activate && pip install -r requirements.txt && deactivate
  python3 scripts/build_bnhub_premium_pitch_deck.py

Apply Fade transitions in PowerPoint manually (python-pptx does not set them).
"""

from __future__ import annotations

import io
import os
import shutil
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
PRES = ROOT / "docs" / "presentations"
SCREEN = PRES / "screenshots"
ARCHIVE = PRES / "archive"
SOURCE_DECK = PRES / "BNHub_Pitch_Deck.pptx"
OUT_DECK = PRES / "BNHub_Pitch_Deck_Premium.pptx"

# After building the template, copy shapes **verbatim** from SOURCE_DECK onto each slide (same positions, images, tables).
# Disable: `PITCH_MERGE=0 python3 scripts/build_bnhub_premium_pitch_deck.py`
MERGE_FROM_SOURCE_CONTENT = os.environ.get("PITCH_MERGE", "1").strip() not in ("0", "false", "no")

# Premium slide index (0-based) -> source slide index (0-based). None = use `default_merge_index_map`.
MERGE_TARGET_TO_SOURCE: dict[int, int] | None = None

BLACK = RGBColor(0x0B, 0x0B, 0x0B)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xBF, 0xBF, 0xBF)
PANEL = RGBColor(0x1A, 0x1A, 0x1A)


def set_slide_bg(slide, color: RGBColor = BLACK) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_block(
    slide,
    title: str,
    bullets: list[str],
    *,
    left_in: float = 0.65,
    top_in: float = 0.55,
    width_in: float = 5.8,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(5.8)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_bottom = Pt(8)
    tf.margin_top = Pt(4)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(32)
    p0.font.bold = True
    p0.font.color.rgb = WHITE
    p0.space_after = Pt(20)

    for line in bullets:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(15)
        p.font.color.rgb = MUTED
        p.space_after = Pt(10)
        p.line_spacing = 1.15


def add_footer(slide, text: str = "Confidential · BNHub + LECIPM") -> None:
    box = slide.shapes.add_textbox(
        Inches(0.65), Inches(7.05), Inches(12), Inches(0.35)
    )
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def ensure_placeholders() -> None:
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return

    specs = [
        ("homepage.png", "Homepage — replace with capture"),
        ("listing.png", "Listing — replace with capture"),
        ("map_search.png", "Map / search — replace with capture"),
        ("dashboard.png", "Dashboard — replace with capture"),
    ]
    SCREEN.mkdir(parents=True, exist_ok=True)
    for fname, label in specs:
        path = SCREEN / fname
        if path.exists():
            continue
        w, h = 1280, 800
        img = Image.new("RGB", (w, h), (11, 11, 11))
        draw = ImageDraw.Draw(img)
        margin = 24
        draw.rectangle(
            [margin, margin, w - margin, h - margin],
            outline=(212, 175, 55),
            width=4,
        )
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Supplemental/Arial.ttf", 28)
        except OSError:
            font = ImageFont.load_default()
        try:
            bbox = draw.textbbox((0, 0), label, font=font)
            tw = bbox[2] - bbox[0]
        except Exception:
            tw = len(label) * 12
        draw.text(((w - tw) // 2, h // 2 - 20), label, fill=(191, 191, 191), font=font)
        img.save(path, "PNG")


def backup_source_if_present() -> None:
    if not SOURCE_DECK.is_file():
        return
    ARCHIVE.mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    dest = ARCHIVE / f"BNHub_Pitch_Deck_backup_{ts}.pptx"
    shutil.copy2(SOURCE_DECK, dest)


def default_merge_index_map(num_dest_slides: int, num_source_slides: int) -> dict[int, int]:
    """
    Map premium slide index -> BNHub_Pitch_Deck slide index.

    Repo source deck has 17 slides (0–16): Vision at 14, The opportunity at 15, Ask at 16.
    Premium template has 16 slides (0–15): single closing vision slide at 14, Ask at 15.
    """
    m: dict[int, int] = {}
    if num_dest_slides == 16 and num_source_slides >= 17:
        for i in range(14):
            m[i] = i
        m[14] = 15  # The opportunity (skip standalone Vision at 14)
        m[15] = 16  # The ask
        return m
    for i in range(min(num_dest_slides, num_source_slides)):
        m[i] = i
    return m


def remove_all_shapes(slide) -> None:
    for shape in reversed(list(slide.shapes)):
        el = shape._element
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)


def _copy_text_frame(source_tf, dest_tf) -> None:
    if hasattr(dest_tf, "clear"):
        dest_tf.clear()
    for i, sp in enumerate(source_tf.paragraphs):
        dp = dest_tf.paragraphs[0] if i == 0 else dest_tf.add_paragraph()
        dp.text = sp.text
        dp.level = sp.level
        try:
            dp.line_spacing = sp.line_spacing
        except Exception:
            pass
        try:
            dp.space_after = sp.space_after
            dp.space_before = sp.space_before
        except Exception:
            pass
        if sp.alignment is not None:
            dp.alignment = sp.alignment
        if sp.runs and dp.runs:
            for sr, dr in zip(sp.runs, dp.runs):
                try:
                    if sr.font.size:
                        dr.font.size = sr.font.size
                except Exception:
                    pass
                try:
                    if sr.font.bold is not None:
                        dr.font.bold = sr.font.bold
                except Exception:
                    pass
                try:
                    rgb = getattr(sr.font.color, "rgb", None)
                    if rgb is not None:
                        dr.font.color.rgb = rgb
                except Exception:
                    pass


def copy_shape_from_source(src_shape, dest_slide) -> None:
    st = src_shape.shape_type
    left, top, w, h = src_shape.left, src_shape.top, src_shape.width, src_shape.height

    if st == MSO_SHAPE_TYPE.PICTURE:
        dest_slide.shapes.add_picture(
            io.BytesIO(src_shape.image.blob), left, top, width=w, height=h
        )
        return

    if getattr(src_shape, "has_table", False):
        tbl = src_shape.table
        rows, cols = len(tbl.rows), len(tbl.columns)
        graphic = dest_slide.shapes.add_table(rows, cols, left, top, w, h)
        nt = graphic.table
        for r in range(rows):
            for c in range(cols):
                try:
                    nt.cell(r, c).text = tbl.cell(r, c).text
                except Exception:
                    pass
        return

    if src_shape.has_text_frame:
        box = dest_slide.shapes.add_textbox(left, top, w, h)
        box.text_frame.word_wrap = src_shape.text_frame.word_wrap
        _copy_text_frame(src_shape.text_frame, box.text_frame)
        return

    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        try:
            ns = dest_slide.shapes.add_shape(
                src_shape.auto_shape_type, left, top, w, h
            )
            if src_shape.fill.type is not None:
                try:
                    if src_shape.fill.type == 1:  # solid
                        ns.fill.solid()
                        fr = getattr(src_shape.fill.fore_color, "rgb", None)
                        if fr is not None:
                            ns.fill.fore_color.rgb = fr
                except Exception:
                    pass
            try:
                lr = getattr(src_shape.line.color, "rgb", None)
                if lr is not None:
                    ns.line.color.rgb = lr
            except Exception:
                pass
            if src_shape.has_text_frame:
                _copy_text_frame(src_shape.text_frame, ns.text_frame)
        except Exception:
            pass
        return


def merge_source_into_premium(dest_prs: Presentation, src_prs: Presentation) -> None:
    n_d, n_s = len(dest_prs.slides), len(src_prs.slides)
    idx_map = MERGE_TARGET_TO_SOURCE or default_merge_index_map(n_d, n_s)
    print(f"Merge: applying source slide map ({len(idx_map)} slides) from {SOURCE_DECK.name}")
    for dest_i in sorted(idx_map.keys()):
        src_i = idx_map[dest_i]
        if dest_i >= n_d or src_i >= n_s:
            print(f"  skip dest {dest_i} <- src {src_i} (out of range)")
            continue
        dest_slide = dest_prs.slides[dest_i]
        src_slide = src_prs.slides[src_i]
        remove_all_shapes(dest_slide)
        for sh in src_slide.shapes:
            try:
                copy_shape_from_source(sh, dest_slide)
            except Exception as exc:
                print(f"  warn: dest {dest_i} <- src {src_i} shape skipped: {exc}")
        set_slide_bg(dest_slide)


def add_table_slide(slide, headers: tuple[str, str], rows: list[tuple[str, str]], title: str) -> None:
    set_slide_bg(slide)
    add_title_block(slide, title, [])
    rows_n = len(rows) + 1
    cols_n = 2
    left, top = Inches(0.65), Inches(1.85)
    width, height = Inches(12), Inches(4.2)
    table = slide.shapes.add_table(rows_n, cols_n, left, top, width, height).table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE if c == 1 else MUTED

    for r, (a, b) in enumerate(rows, start=1):
        table.cell(r, 0).text = a
        table.cell(r, 1).text = b
        for c, txt in enumerate((a, b)):
            cell = table.cell(r, c)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = MUTED if c == 0 else WHITE

    add_footer(slide)


def build() -> None:
    backup_source_if_present()
    ensure_placeholders()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # --- 1 Title ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    tb = s.shapes.add_textbox(Inches(1), Inches(1.9), Inches(11.3), Inches(3.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "BNHub"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p2 = tf.add_paragraph()
    p2.text = "+ LECIPM"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(28)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(8)
    p3 = tf.add_paragraph()
    p3.text = "AI-powered real estate & rental platform"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(17)
    p3.font.color.rgb = MUTED
    p3.space_before = Pt(20)
    p4 = tf.add_paragraph()
    p4.text = "[Founder name] · [Company] · Confidential"
    p4.alignment = PP_ALIGN.CENTER
    p4.font.size = Pt(12)
    p4.font.color.rgb = MUTED
    p4.space_before = Pt(28)
    logo_note = s.shapes.add_textbox(Inches(1), Inches(0.85), Inches(11.3), Inches(0.4))
    logo_note.text_frame.paragraphs[0].text = "[BNHub logo placeholder]"
    logo_note.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    logo_note.text_frame.paragraphs[0].font.size = Pt(11)
    logo_note.text_frame.paragraphs[0].font.color.rgb = GOLD

    # --- 2 Problem ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_block(
        s,
        "Fragmentation erodes trust",
        [
            "Listings, bookings, and pro tools in silos",
            "Guests: unclear fees, slow answers",
            "Operators: manual work that does not scale",
        ],
    )
    # Right: simple diagram — 3 boxes labels
    for i, lab in enumerate(["Listings", "Bookings", "Ops"]):
        y = 1.4 + i * 1.45
        sh = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.1),
            Inches(y),
            Inches(5.2),
            Inches(0.85),
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = PANEL
        sh.line.color.rgb = GOLD
        tf = sh.text_frame
        tf.paragraphs[0].text = lab
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_footer(s)

    # --- 3 Solution ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_block(
        s,
        "One platform: property + stays",
        [
            "Unified real estate + short-term rentals",
            "Clearer pricing and fewer handoffs",
            "Depth in priority markets first",
        ],
    )
    sh = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.8), Inches(2.4), Inches(1.1)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = PANEL
    sh.line.color.rgb = GOLD
    sh.text_frame.paragraphs[0].text = "LECIPM"
    sh.text_frame.paragraphs[0].font.color.rgb = WHITE
    sh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    sh2 = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.0), Inches(1.8), Inches(2.4), Inches(1.1)
    )
    sh2.fill.solid()
    sh2.fill.fore_color.rgb = PANEL
    sh2.line.color.rgb = GOLD
    sh2.text_frame.paragraphs[0].text = "BNHub"
    sh2.text_frame.paragraphs[0].font.color.rgb = WHITE
    sh2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    bridge = s.shapes.add_textbox(Inches(9.35), Inches(2.15), Inches(1), Inches(0.5))
    bridge.text_frame.paragraphs[0].text = "↔"
    bridge.text_frame.paragraphs[0].font.size = Pt(28)
    bridge.text_frame.paragraphs[0].font.color.rgb = GOLD
    add_footer(s)

    # --- 4 Product (screenshots) ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_block(
        s,
        "Product",
        [
            "BNHub: search, listing, booking, host tools",
            "LECIPM: listings, broker & owner workflows",
        ],
    )
    paths = [
        (SCREEN / "homepage.png", 6.85, 1.35, 2.95, 2.05),
        (SCREEN / "listing.png", 9.95, 1.35, 2.95, 2.05),
        (SCREEN / "map_search.png", 6.85, 3.55, 2.95, 2.05),
        (SCREEN / "dashboard.png", 9.95, 3.55, 2.95, 2.05),
    ]
    for path, lx, ty, w, h in paths:
        if path.is_file():
            pl, pt, pw, ph = Inches(lx), Inches(ty), Inches(w), Inches(h)
            fr = s.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                pl - Inches(0.05),
                pt - Inches(0.05),
                pw + Inches(0.1),
                ph + Inches(0.1),
            )
            fr.fill.solid()
            fr.fill.fore_color.rgb = PANEL
            fr.line.color.rgb = GOLD
            fr.line.width = Pt(1)
            s.shapes.add_picture(str(path), pl, pt, width=pw, height=ph)
    add_footer(s)

    # --- 5 How it works ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_block(
        s,
        "How it works",
        [],
    )
    steps = [
        ("1", "List property", "Supply on LECIPM / BNHub"),
        ("2", "Discover & book", "Guest search and checkout"),
        ("3", "Manage & optimize", "Calendar, payouts, pro tools"),
    ]
    for i, (num, t, sub) in enumerate(steps):
        x = 0.85 + i * 4.0
        circ = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(2.0), Inches(0.75), Inches(0.75)
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = GOLD
        circ.line.fill.background()
        circ.text_frame.paragraphs[0].text = num
        circ.text_frame.paragraphs[0].font.size = Pt(18)
        circ.text_frame.paragraphs[0].font.bold = True
        circ.text_frame.paragraphs[0].font.color.rgb = BLACK
        circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb = s.shapes.add_textbox(Inches(x - 0.2), Inches(2.95), Inches(3.6), Inches(1.8))
        tf = tb.text_frame
        tf.paragraphs[0].text = t
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(13)
        p2.font.color.rgb = MUTED
    add_footer(s)

    # --- 6 Market (big numbers) ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_block(
        s,
        "Market opportunity",
        [
            "Large STR + proptech categories (directional)",
            "Growing demand for flexible stays",
            "We win with bottom-up city depth",
        ],
    )
    metrics = [("TAM", "$XB+", "illustrative"), ("SAM", "$XM", "focus geography"), ("Growth", "DD%", "public sources")]
    for i, (lab, num, sub) in enumerate(metrics):
        bx = s.shapes.add_textbox(Inches(6.8 + i * 2.05), Inches(1.9), Inches(1.85), Inches(1.6))
        tf = bx.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(26)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = GOLD
        p2 = tf.add_paragraph()
        p2.text = lab
        p2.font.size = Pt(12)
        p2.font.color.rgb = WHITE
        p3 = tf.add_paragraph()
        p3.text = sub
        p3.font.size = Pt(10)
        p3.font.color.rgb = MUTED
    # simple bar visual
    bar_y = 4.0
    for i, w in enumerate([4.5, 3.2, 1.8]):
        bar = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.8),
            Inches(bar_y + i * 0.55),
            Inches(w),
            Inches(0.35),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = GOLD if i == 2 else PANEL
        bar.line.color.rgb = GOLD
    add_footer(s)

    # --- 7 Business model ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_block(
        s,
        "Business model",
        [
            "Booking commissions",
            "Premium host & pro features",
            "Featured placement when ROI-positive",
        ],
    )
    add_footer(s)

    # --- 8 Traction ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_block(
        s,
        "Traction",
        [
            "Early validation phase — replace with facts",
            "Users / listings: [as of date]",
            "Next milestone: [90-day goal]",
        ],
    )
    call = s.shapes.add_textbox(Inches(6.8), Inches(1.85), Inches(5.8), Inches(2.2))
    tf = call.text_frame
    tf.paragraphs[0].text = "Early validation"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GOLD
    p2 = tf.add_paragraph()
    p2.text = "Tests · onboarding · pilot hosts"
    p2.font.size = Pt(15)
    p2.font.color.rgb = MUTED
    add_footer(s)

    # --- 9 GTM ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_block(
        s,
        "Go-to-market",
        [
            "Montreal first — density & relationships",
            "Direct onboarding + content + referrals",
            "Prove repeat bookings, then repeat playbook",
        ],
    )
    add_footer(s)

    # --- 10 Competition ---
    s = prs.slides.add_slide(blank)
    add_table_slide(
        s,
        ("Typical OTA (e.g. Airbnb)", "BNHub + LECIPM"),
        [
            ("Booking-first consumer habit", "Ecosystem: RE + STR integrated"),
            ("Limited broker / owner tooling", "Pro workflows + stays in one stack"),
            ("Global scale; less local ops depth", "Local execution, partner GTM"),
        ],
        "Competitive positioning",
    )

    # --- 11 Technology ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_block(
        s,
        "Technology",
        [
            "Unified codebase & APIs",
            "Web-first; mobile per roadmap",
            "AI where shipped — labeled vs roadmap",
        ],
    )
    add_footer(s)

    # --- 12 Growth ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_block(
        s,
        "Growth strategy",
        [
            "Same playbook: prove → document → expand",
            "B2B / agency partnerships",
            "Cut channels that do not move bookings",
        ],
    )
    add_footer(s)

    # --- 13 Financials flow ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_block(
        s,
        "Financials (high level)",
        [
            "All forward figures = estimates — see appendix",
        ],
    )
    flow = [("Users", 0.5), ("Bookings", 4.5), ("Commission", 8.5)]
    for i, (lab, x) in enumerate(flow):
        bx = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(2.4),
            Inches(2.8),
            Inches(0.9),
        )
        bx.fill.solid()
        bx.fill.fore_color.rgb = PANEL
        bx.line.color.rgb = GOLD
        bx.text_frame.paragraphs[0].text = lab
        bx.text_frame.paragraphs[0].font.size = Pt(15)
        bx.text_frame.paragraphs[0].font.color.rgb = WHITE
        bx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        if i < len(flow) - 1:
            ar = s.shapes.add_textbox(Inches(x + 2.95), Inches(2.65), Inches(0.8), Inches(0.5))
            ar.text_frame.paragraphs[0].text = "→"
            ar.text_frame.paragraphs[0].font.size = Pt(28)
            ar.text_frame.paragraphs[0].font.color.rgb = GOLD
    add_footer(s)

    # --- 14 Team ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_block(
        s,
        "Team",
        [
            "Founder: licensed real estate broker",
            "Domain depth on closings, listings, clients",
            "Hiring: eng, GTM, trust ops",
        ],
    )
    add_footer(s)

    # --- 15 The Opportunity ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    tb = s.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(10.9), Inches(2.8))
    tf = tb.text_frame
    tf.paragraphs[0].text = "The Opportunity"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = GOLD
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = (
        "Connect real estate and hospitality\nin one trusted ecosystem."
    )
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    add_footer(s)

    # --- 16 Ask ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_title_block(
        s,
        "The ask",
        [
            "Raising: [amount] — [structure]",
            "Product & engineering",
            "Growth & Montreal execution",
            "[email] · [calendar]",
        ],
    )
    add_footer(s)

    PRES.mkdir(parents=True, exist_ok=True)
    if MERGE_FROM_SOURCE_CONTENT and SOURCE_DECK.is_file():
        src_prs = Presentation(str(SOURCE_DECK))
        merge_source_into_premium(prs, src_prs)
    prs.save(str(OUT_DECK))
    print(f"Wrote {OUT_DECK}")


if __name__ == "__main__":
    build()
