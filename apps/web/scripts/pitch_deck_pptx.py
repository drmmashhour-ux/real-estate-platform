#!/usr/bin/env python3
"""
Optional PPTX export using python-pptx (not used by Next.js runtime).

  pip install python-pptx
  python scripts/pitch_deck_pptx.py --input deck.json --output /mnt/data/lecipm_pitch_deck.pptx

`deck.json` should be `{ "slides": [ { "title", "type", "order", "content": { "bullets": [] } } ] }`
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path


def main() -> None:
    parser = argparse.ArgumentParser(description="Build LECIPM pitch deck PPTX (python-pptx).")
    parser.add_argument("--input", required=True, help="JSON file with slides array")
    parser.add_argument(
        "--output",
        default="/mnt/data/lecipm_pitch_deck.pptx",
        help="Output path (default matches Colab / Cursor mount)",
    )
    args = parser.parse_args()

    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.util import Inches, Pt  # type: ignore[import-untyped]
    except ImportError:
        print("Install python-pptx: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    raw = Path(args.input).read_text(encoding="utf-8")
    data = json.loads(raw)
    slides = sorted(data.get("slides", []), key=lambda s: s.get("order", 0))

    prs = Presentation()
    prs.core_properties.title = "LECIPM Investor Pitch"
    prs.core_properties.author = "LECIPM"

    for spec in slides:
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(1))
        tf = title_box.text_frame
        tf.text = str(spec.get("title", ""))
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True

        content = spec.get("content") or {}
        lines: list[str] = []
        if isinstance(content, dict):
            if isinstance(content.get("subtitle"), str):
                lines.append(content["subtitle"])
            bullets = content.get("bullets")
            if isinstance(bullets, list):
                if lines:
                    lines.append("")
                lines.extend(f"• {b}" for b in bullets if isinstance(b, str))

        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(12), Inches(5.5))
        btf = body.text_frame
        btf.text = "\n".join(lines) if lines else json.dumps(content, indent=2)
        for p in btf.paragraphs:
            p.font.size = Pt(16)

    out = Path(args.output)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
