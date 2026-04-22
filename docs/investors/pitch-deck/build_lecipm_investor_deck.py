#!/usr/bin/env python3
"""
Build LECIPM_Investor_Deck.pptx — 1920×1080 equivalent (16:9), black + gold brand.

Prereq (from repo root):
  pip install -r docs/investors/pitch-deck/requirements.txt

Run:
  python docs/investors/pitch-deck/build_lecipm_investor_deck.py
  python docs/investors/pitch-deck/build_lecipm_investor_deck.py --output ./LECIPM_Investor_Deck.pptx

Fonts: Install Montserrat + Open Sans on the machine opening the file; fallback is Calibri.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("Missing python-pptx. Run: pip install -r docs/investors/pitch-deck/requirements.txt", file=sys.stderr)
    sys.exit(1)

# Brand
BLACK = RGBColor(0x00, 0x00, 0x00)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
GOLD2 = RGBColor(0xC9, 0xA6, 0x46)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUBTEXT = RGBColor(0xBF, 0xBF, 0xBF)
PANEL = RGBColor(0x11, 0x11, 0x11)

# 1920×1080 @ 13.333×7.5 in mapping
SW = Inches(13.3333333)
SH = Inches(7.5)


def px_w(px: float) -> float:
    return float(px) / 1920.0 * 13.3333333


def px_h(px: float) -> float:
    return float(px) / 1080.0 * 7.5


M_LEFT = px_w(80)
M_TOP = px_h(80)
M_RIGHT = px_w(80)
BODY = Pt(22)
TITLE = Pt(52)
TITLE_HERO = Pt(54)
SUB = Pt(28)
LABEL = Pt(14)
SMALL = Pt(18)
KPI_NUM = Pt(60)


def set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BLACK


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def style_run(run, *, size=BODY, bold=False, color=WHITE, name="Open Sans"):
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_footer(slide, text: str = "Confidential · LECIPM") -> None:
    box = slide.shapes.add_textbox(M_LEFT, SH - px_h(70), SW - M_LEFT - M_RIGHT, px_h(28))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    style_run(p.runs[0], size=Pt(11), bold=False, color=SUBTEXT)


def slide_title_block(slide, label: str, title: str, *, top_off: float = 0) -> None:
    y = M_TOP + px_h(top_off)
    if label:
        lb = slide.shapes.add_textbox(M_LEFT, y, SW - M_LEFT - M_RIGHT, px_h(24))
        tf = lb.text_frame
        p = tf.paragraphs[0]
        p.text = label.upper()
        style_run(p.runs[0], size=LABEL, bold=True, color=GOLD2, name="Montserrat")
        p.space_after = Pt(8)
        y += px_h(36)

    if title and title.strip():
        tb = slide.shapes.add_textbox(M_LEFT, y, SW - M_LEFT - M_RIGHT - px_w(560), px_h(120))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        style_run(p.runs[0], size=TITLE, bold=True, color=WHITE, name="Montserrat")
        p.space_after = Pt(16)


def slide_1(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide)
    box = slide.shapes.add_textbox(M_LEFT, px_h(340), SW - M_LEFT - M_RIGHT, px_h(200))
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "LECIPM"
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], size=TITLE_HERO, bold=True, color=WHITE, name="Montserrat")

    sub = slide.shapes.add_textbox(M_LEFT, px_h(460), SW - M_LEFT - M_RIGHT, px_h(80))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "AI-Powered Real Estate Operating System"
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], size=SUB, bold=True, color=SUBTEXT, name="Montserrat")

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        SW / 2 - px_w(100),
        px_h(540),
        px_w(200),
        px_h(2),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    tag = slide.shapes.add_textbox(M_LEFT, px_h(570), SW - M_LEFT - M_RIGHT, px_h(40))
    tp = tag.text_frame.paragraphs[0]
    tp.text = "Confidential · Investor presentation"
    tp.alignment = PP_ALIGN.CENTER
    style_run(tp.runs[0], size=SMALL, bold=False, color=SUBTEXT)


def slide_2(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide)
    slide_title_block(slide, "Problem", "The stack stopped at listings")

    bullets = [
        "Broken system: passive directories don’t execute deals.",
        "Lost leads: attribution drops after the first click.",
        "Inefficiency: brokers rely on habit, not leverage.",
        "No unified intelligence across CRM, deals, marketplace.",
        "Regulators expect auditability — spreadsheets don’t scale.",
    ]
    bx = slide.shapes.add_textbox(M_LEFT, px_h(220), px_w(900), px_h(520))
    tf = bx.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.space_after = Pt(14)
        p.line_spacing = 1.35
        style_run(p.runs[0], size=BODY, bold=False, color=WHITE)

    captions = ["Broken\nsystem", "Lost\nleads", "Inefficiency"]
    x0 = px_w(1120)
    for i, cap in enumerate(captions):
        iy = px_h(260 + i * 140)
        r = slide.shapes.add_shape(MSO_SHAPE.OVAL, x0 + px_w(60), iy, px_w(100), px_w(100))
        r.fill.background()
        r.line.color.rgb = GOLD
        r.line.width = Pt(2)

        cb = slide.shapes.add_textbox(x0, iy + px_w(110), px_w(200), px_h(80))
        cp = cb.text_frame.paragraphs[0]
        cp.text = cap.replace("\n", " ")
        cp.alignment = PP_ALIGN.CENTER
        style_run(cp.runs[0], size=Pt(18), bold=False, color=SUBTEXT)


def slide_3(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide)
    slide_title_block(slide, "Market gap", "Distribution ≠ execution")

    cols = [
        ("Centris", "Reach and syndication. Execution still happens elsewhere."),
        ("CRM", "Storage and tasks. Limited intelligence."),
        ("Missing", "Intelligence layer — conversion, workflow depth, learning.", True),
    ]
    gap = px_w(40)
    cw = (SW - M_LEFT - M_RIGHT - gap * 2) / 3
    y = px_h(260)
    h = px_h(340)
    for i, (head, body, *rest) in enumerate(cols):
        hi = rest[0] if rest else False
        left = M_LEFT + i * (cw + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, cw, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = PANEL
        shp.line.color.rgb = GOLD if hi else RGBColor(0x55, 0x55, 0x55)
        shp.line.width = Pt(2 if hi else 1)

        tb = slide.shapes.add_textbox(left + px_w(24), y + px_h(28), cw - px_w(48), h - px_h(40))
        tf = tb.text_frame
        p0 = tf.paragraphs[0]
        p0.text = head
        style_run(p0.runs[0], size=Pt(28), bold=True, color=GOLD if hi else WHITE, name="Montserrat")
        p1 = tf.add_paragraph()
        p1.text = body
        p1.space_before = Pt(16)
        style_run(p1.runs[0], size=Pt(20), bold=False, color=SUBTEXT)


def slide_4(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide)
    slide_title_block(slide, "Solution", "One operating layer for transactions")

    bullets = [
        "Captures & attributes leads end-to-end.",
        "Analyzes deals — scores, risk, next steps (assistive AI).",
        "Guides brokers — suggestions, not autonomous outreach.",
        "Learns from outcomes — tighter defaults over time.",
    ]
    bx = slide.shapes.add_textbox(M_LEFT, px_h(220), SW - M_LEFT - M_RIGHT, px_h(520))
    tf = bx.text_frame
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"●  {line}"
        p.space_after = Pt(22)
        p.line_spacing = 1.4
        style_run(p.runs[0], size=BODY, bold=False, color=WHITE)
        p.runs[0].font.color.rgb = GOLD
        # second run for body after bullet — simplify: one run white
        p.clear()
        p.text = ""
        r1 = p.add_run()
        r1.text = "●  "
        style_run(r1, size=BODY, bold=True, color=GOLD)
        r2 = p.add_run()
        r2.text = line
        style_run(r2, size=BODY, bold=False, color=WHITE)


def add_flow(slide, nodes: list[str], *, y: float, label_above: bool = False):
    total_w = SW - M_LEFT - M_RIGHT - px_w(100)
    n = len(nodes)
    slot = total_w / max(n, 1)
    x = M_LEFT + px_w(40)
    for i, node in enumerate(nodes):
        cx = x + i * slot + slot / 2 - px_w(70)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y, px_w(140), px_h(52))
        box.adjustments[0] = 0.08
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
        box.line.color.rgb = GOLD
        tf = box.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = node
        style_run(p.runs[0], size=Pt(18), bold=True, color=WHITE, name="Montserrat")

        if i < n - 1:
            ax = cx + px_w(150)
            ab = slide.shapes.add_textbox(ax, y + px_h(12), px_w(36), px_h(40))
            ap = ab.text_frame.paragraphs[0]
            ap.text = "→"
            style_run(ap.runs[0], size=Pt(22), bold=True, color=GOLD2, name="Montserrat")


def slide_5(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide)
    slide_title_block(slide, "Product", "Overview")
    add_flow(slide, ["Lead", "Deal", "AI", "Broker", "Close"], y=px_h(400))

    cap = slide.shapes.add_textbox(M_LEFT, px_h(500), SW - M_LEFT - M_RIGHT, px_h(80))
    cp = cap.text_frame.paragraphs[0]
    cp.text = "Modular core: funnel → intelligence → assistant → compliance → learning."
    cp.alignment = PP_ALIGN.CENTER
    style_run(cp.runs[0], size=Pt(20), bold=False, color=SUBTEXT)


def slide_6(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide)
    slide_title_block(slide, "How it works", "From syndication to closed deal")
    add_flow(
        slide,
        ["Centris", "LECIPM", "AI Engine", "Broker", "Closed Deal"],
        y=px_h(400),
    )
    cap = slide.shapes.add_textbox(M_LEFT, px_h(500), SW - M_LEFT - M_RIGHT, px_h(60))
    cp = cap.text_frame.paragraphs[0]
    cp.text = "Partner data only as permitted by contract."
    cp.alignment = PP_ALIGN.CENTER
    style_run(cp.runs[0], size=Pt(18), bold=False, color=SUBTEXT)


def slide_7(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide)
    slide_title_block(slide, "Differentiation", "Platforms vs. operating system")
    y = px_h(220)
    h = px_h(460)
    w = (SW - M_LEFT - M_RIGHT - px_w(40)) / 2

    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, M_LEFT, y, w, h)
    left.fill.solid()
    left.fill.fore_color.rgb = RGBColor(0x0C, 0x0C, 0x0C)
    left.line.color.rgb = RGBColor(0x55, 0x55, 0x55)
    tf = slide.shapes.add_textbox(M_LEFT + px_w(28), y + px_h(28), w - px_w(56), h - px_h(56)).text_frame
    p = tf.paragraphs[0]
    p.text = "Traditional platforms"
    style_run(p.runs[0], size=Pt(24), bold=True, color=SUBTEXT, name="Montserrat")
    for line in [
        "— Listings-first or CRM-only",
        "— Weak cross-object intelligence",
        "— Limited outcome learning",
    ]:
        pr = tf.add_paragraph()
        pr.text = line
        pr.space_before = Pt(12)
        style_run(pr.runs[0], size=Pt(20), bold=False, color=SUBTEXT)

    rx = M_LEFT + w + px_w(40)
    right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rx, y, w, h)
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(0x12, 0x10, 0x08)
    right.line.color.rgb = GOLD
    right.line.width = Pt(2)
    tf2 = slide.shapes.add_textbox(rx + px_w(28), y + px_h(28), w - px_w(56), h - px_h(56)).text_frame
    p = tf2.paragraphs[0]
    p.text = "LECIPM advantages"
    style_run(p.runs[0], size=Pt(24), bold=True, color=GOLD, name="Montserrat")
    for line in [
        "✓ Transaction OS — leads through close",
        "✓ Learning + marketplace signals",
        "✓ Québec-aligned compliance posture",
        "✓ Daily broker workflow depth",
    ]:
        pr = tf2.add_paragraph()
        pr.text = line
        pr.space_before = Pt(12)
        style_run(pr.runs[0], size=Pt(20), bold=False, color=WHITE)


def slide_8(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide)
    slide_title_block(slide, "Business model", "Revenue architecture")
    cards = [
        ("Subscription", "Broker seats and workspace."),
        ("AI tools", "Premium scoring, exports, assistant depth."),
        ("Analytics", "Investor surfaces (non-advisory)."),
        ("Transaction fees", "Aligned to closed flow where permitted."),
    ]
    gw = px_w(40)
    cw = (SW - M_LEFT - M_RIGHT - gw) / 2
    ch = px_h(200)
    y0 = px_h(240)
    for idx, (title, desc) in enumerate(cards):
        row, col = divmod(idx, 2)
        left = M_LEFT + col * (cw + gw)
        top = y0 + row * (ch + px_h(40))
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, cw, ch)
        r.fill.solid()
        r.fill.fore_color.rgb = PANEL
        r.line.color.rgb = GOLD
        tf = slide.shapes.add_textbox(left + px_w(28), top + px_h(24), cw - px_w(56), ch - px_h(40)).text_frame
        p0 = tf.paragraphs[0]
        p0.text = title
        style_run(p0.runs[0], size=Pt(24), bold=True, color=GOLD, name="Montserrat")
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.space_before = Pt(10)
        style_run(p1.runs[0], size=Pt(20), bold=False, color=SUBTEXT)


def slide_9(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide)
    slide_title_block(slide, "Traction", "Proof points")
    labels = [(" [#] ", "LEADS"), (" [#] ", "BROKERS"), (" [#] ", "DEALS")]
    slot = (SW - M_LEFT - M_RIGHT) / 3
    y = px_h(360)
    for i, (num, lab) in enumerate(labels):
        left = M_LEFT + i * slot
        tb = slide.shapes.add_textbox(left, y, slot - px_w(20), px_h(140))
        tf = tb.text_frame
        p0 = tf.paragraphs[0]
        p0.text = num.strip()
        p0.alignment = PP_ALIGN.LEFT
        style_run(p0.runs[0], size=KPI_NUM, bold=True, color=GOLD, name="Montserrat")
        p1 = tf.add_paragraph()
        p1.text = lab
        p1.space_before = Pt(14)
        style_run(p1.runs[0], size=Pt(20), bold=True, color=SUBTEXT)


def slide_10(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide)
    slide_title_block(slide, "Moat", "Why this compounds")
    items = [
        ("Data ", "Integrated graph — leads ↔ deals ↔ listings."),
        ("AI ", "Self-improving signals with human review."),
        ("Compliance ", "Built for auditability & Québec."),
        ("Workflow ", "Deep daily integration."),
    ]
    y = px_h(230)
    for head, body in items:
        bx = slide.shapes.add_textbox(M_LEFT, y, px_w(900), px_h(72))
        tf = bx.text_frame
        p = tf.paragraphs[0]
        r0 = p.add_run()
        r0.text = head
        style_run(r0, size=BODY, bold=True, color=GOLD, name="Montserrat")
        r1 = p.add_run()
        r1.text = body
        style_run(r1, size=BODY, bold=False, color=WHITE)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, M_LEFT, y + px_h(52), px_w(880), px_h(1))
        line.fill.solid()
        line.fill.fore_color.rgb = GOLD
        line.line.fill.background()
        y += px_h(88)


def slide_11(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide)
    slide_title_block(slide, "Vision", " ")
    # Clear subtitle hack: empty title area — fix by second textbox only

    lines = [
        "Depth in priority markets first — prove the OS where it matters.",
        "Assistive AI scales; licensees retain decision rights.",
        "Expand modules where product and regulation align.",
    ]
    y = px_h(300)
    for line in lines:
        bx = slide.shapes.add_textbox(M_LEFT, y, SW - M_LEFT - M_RIGHT - px_w(100), px_h(100))
        tf = bx.text_frame
        p = tf.paragraphs[0]
        p.text = line
        style_run(p.runs[0], size=Pt(30), bold=True, color=WHITE, name="Montserrat")
        p.line_spacing = 1.35
        y += px_h(110)


def slide_12(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide)
    slide_title_block(slide, "Ask", "The round")
    w = (SW - M_LEFT - M_RIGHT - px_w(60)) / 2
    y = px_h(220)
    h = px_h(420)

    lb = slide.shapes.add_textbox(M_LEFT, y, w, h)
    tf = lb.text_frame
    p0 = tf.paragraphs[0]
    p0.text = "$ [___]"
    style_run(p0.runs[0], size=Pt(48), bold=True, color=GOLD, name="Montserrat")
    p1 = tf.add_paragraph()
    p1.text = "Stage: [Seed / Series A]"
    p1.space_before = Pt(20)
    style_run(p1.runs[0], size=BODY, bold=False, color=SUBTEXT)
    p2 = tf.add_paragraph()
    p2.text = "Product velocity · broker GTM · key hires."
    p2.space_before = Pt(12)
    style_run(p2.runs[0], size=BODY, bold=False, color=WHITE)

    rx = M_LEFT + w + px_w(60)
    pie = slide.shapes.add_shape(MSO_SHAPE.PIE, rx + px_w(40), y + px_h(60), px_w(200), px_h(200))
    pie.fill.solid()
    pie.fill.fore_color.rgb = GOLD
    pie.line.color.rgb = GOLD2

    legend = slide.shapes.add_textbox(rx + px_w(280), y + px_h(80), w - px_w(300), h)
    tf = legend.text_frame
    rows = [
        "Product · 45%",
        "GTM · 30%",
        "Team · 15%",
        "Buffer · 10%",
    ]
    for i, row in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = row
        p.space_after = Pt(12)
        style_run(p.runs[0], size=Pt(20), bold=False, color=WHITE)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("docs/investors/pitch-deck/LECIPM_Investor_Deck.pptx"),
        help="Output .pptx path",
    )
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    prs.core_properties.title = "LECIPM Investor Pitch"
    prs.core_properties.author = "LECIPM"

    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)
    slide_7(prs)
    slide_8(prs)
    slide_9(prs)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs)

    out = args.output.resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
